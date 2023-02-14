import pandas as pd
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN

# Light Style 1 - Accent 6
TABLE_STYLE = "{68D230F3-CF80-4859-8CE7-A43EE81993B5}"


def set_graph(slide, img_path: str):
    """
    Add a graph to a slides PicturePlaceholder.

    Parameters
    ----------
    slide : the slide where the table is to be added to
    img_path : the graphs location

    Returns
    -------
    None
    """
    for shape in slide.shapes:
        if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            shape.insert_picture(img_path)


def set_table(slide, additional: dict):
    """
    Add a table with the contents of additional to a slides TablePlaceholder.

    Parameters
    ----------
    slide : the slide where the table is to be added to
    additional : the contents of the new table.

    Returns
    -------
    None
    """
    for shape in slide.shapes:
        if shape.placeholder_format.type == PP_PLACEHOLDER.TABLE:
            df = pd.DataFrame(additional)
            graphic_frame = shape.insert_table(
                rows=1 + len(df.axes[0]), cols=len(df.axes[1])
            )

            # set table style
            tbl = graphic_frame._element.graphic.graphicData.tbl
            tbl[0][-1].text = TABLE_STYLE

            table = graphic_frame.table

            flat_df = list(df)
            flat_df.extend(df.values.flatten().tolist())

            i = 0
            for cell in table.iter_cells():
                cell.text = str(flat_df[i])
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
                i = i + 1
